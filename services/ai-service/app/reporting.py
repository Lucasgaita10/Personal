"""Report generation: HTML→PDF for memos, python-pptx for decks.

Both formats are branded with the Stone Gate logo (white script on maroon).
The logo lives at app/assets/stone-gate-logo.png and is embedded base64 in
HTML and as a picture in PPTX.
"""
from __future__ import annotations

import base64
import html as html_escape
import os
import re
import uuid
from datetime import datetime
from pathlib import Path

from jinja2 import Environment, BaseLoader
from markupsafe import Markup
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor


# ─── Markdown → HTML ────────────────────────────────────────────────
# The agents (especially market_researcher) emit prose with markdown.
# Without conversion, headings/bold/italic/links render as raw text in
# the PDF. We escape first, then transform a focused set of patterns —
# enough for the agent outputs without a full Markdown library.

def _strip_ai_preamble(text: str) -> str:
    """Drop any 'I'll research…' / 'Let me now compose…' narration before
    the first heading. Same logic as the MarketResearchSection in the UI."""
    m = re.search(r"(?:^|\n)#{1,3}\s+", text)
    if m and m.start() > 0:
        return text[m.start():].lstrip("\n")
    return text


def _render_inline(text: str) -> str:
    """Convert inline markdown patterns to HTML. Operates on already-escaped
    text — placeholders are stripped of escaping so anchors render correctly.
    Order matters: links before bare URLs, bold before italic."""
    # Markdown link [label](url) — match BEFORE bare URL pass to claim the URL.
    text = re.sub(
        r"\[([^\]]+)\]\((https?://[^)\s]+)\)",
        lambda m: f'<a href="{m.group(2)}">{m.group(1)}</a>',
        text,
    )
    # Bare URLs — wrap as anchors. Skip ones already inside href="..." (already linkified above).
    text = re.sub(
        r"(?<!href=\")(?<!\">)(https?://[^\s<)\]]+)",
        r'<a href="\1">\1</a>',
        text,
    )
    # **bold** — non-greedy
    text = re.sub(r"\*\*([^*]+?)\*\*", r"<strong>\1</strong>", text)
    # *italic* — non-greedy, must not start with whitespace
    text = re.sub(r"\*([^*\s][^*]*?)\*", r"<em>\1</em>", text)
    return text


def md_to_html(text: str | None) -> Markup:
    """Convert agent markdown prose to safe HTML for embedding in templates."""
    if not text:
        return Markup("")
    cleaned = _strip_ai_preamble(text)

    # Drop any leading `# Title` heading — the template provides its own.
    cleaned = re.sub(r"^#\s+[^\n]*\n+", "", cleaned)

    blocks: list[str] = []
    # Split on blank lines to get paragraphs / heading blocks
    for raw in re.split(r"\n\s*\n", cleaned.strip()):
        block = raw.strip()
        if not block:
            continue
        # Detect headings on the first line
        h_match = re.match(r"^(#{1,3})\s+(.+?)(\n([\s\S]*))?$", block)
        if h_match:
            level = len(h_match.group(1))
            title = html_escape.escape(h_match.group(2).strip(), quote=False)
            title = _render_inline(title)
            tag = f"h{min(level + 1, 4)}"  # # → h2, ## → h3, ### → h4
            blocks.append(f"<{tag}>{title}</{tag}>")
            rest = (h_match.group(4) or "").strip()
            if rest:
                blocks.append(
                    f"<p>{_render_inline(html_escape.escape(rest, quote=False))}</p>"
                )
        else:
            escaped = html_escape.escape(block, quote=False)
            # Preserve single newlines as <br> for readability
            escaped = escaped.replace("\n", "<br>")
            blocks.append(f"<p>{_render_inline(escaped)}</p>")

    return Markup("\n".join(blocks))

PRIMARY = "#780000"
PRIMARY_RGB = RGBColor(0x78, 0x00, 0x00)
WHITE_RGB = RGBColor(0xFF, 0xFF, 0xFF)

ASSETS_DIR = Path(__file__).parent / "assets"
LOGO_PATH = ASSETS_DIR / "stone-gate-logo.png"
LOGO_WHITE_PATH = ASSETS_DIR / "stone-gate-logo-white.png"


def _logo_data_uri(path: Path = LOGO_PATH) -> str:
    """Read a logo PNG and return a base64 data URI for embedding in HTML."""
    if not path.exists():
        return ""
    raw = path.read_bytes()
    b64 = base64.b64encode(raw).decode("ascii")
    return f"data:image/png;base64,{b64}"


IC_MEMO_TEMPLATE = """\
<!doctype html>
<html><head><meta charset='utf-8'><style>
  /* Zero horizontal page margin so the running header can span edge-to-edge.
     Body adds its own 18mm side padding to keep content indented. */
  @page { size: A4; margin: 32mm 0 22mm 0;
          @top-left-corner { content: ""; }
          @top-left { content: element(header); width: 210mm; padding: 0; margin: 0; }
          @top-right-corner { content: ""; }
          @bottom-right { content: "Page " counter(page) " of " counter(pages);
                          color:#6b6b6b; font-size:8pt; margin-right: 18mm; } }
  body { font-family: 'Inter', 'Helvetica Neue', sans-serif; color:#0a0a0a;
         font-size:10pt; line-height:1.55;
         padding: 0 18mm; }
  h1, h2, h3, h4 { color:#0a0a0a; }

  /* Brand band — repeated on every page via running header.
     Light band with a maroon bottom border, full-width. Logo and label
     sit together on the LEFT (separated by a small gap), vertically
     centered with the logo. */
  .brand-band { position: running(header); background:#ffffff;
                border-bottom: 1.5pt solid #780000;
                padding: 10pt 18mm;
                display:flex; align-items:center; justify-content:flex-start;
                gap: 14pt; }
  .brand-logo { height: 44pt; width: auto; }
  .brand-band .label { font-weight:600; letter-spacing:0.08em; text-transform:uppercase;
                       font-size:9.5pt; color:#780000; }

  h1 { font-size:20pt; margin-top:6pt; margin-bottom:4pt; font-weight:700; }
  h2 { font-size:13pt; border-bottom:1.5px solid #780000; padding-bottom:4pt;
       margin-top:22pt; margin-bottom:10pt; font-weight:600; letter-spacing:0.01em; }
  h3 { font-size:11.5pt; color:#780000; margin-top:14pt; margin-bottom:6pt; font-weight:600; }
  h4 { font-size:10.5pt; color:#0a0a0a; margin-top:10pt; margin-bottom:4pt; font-weight:600; }

  /* Markdown-rendered prose blocks */
  .md p { margin: 0 0 9pt 0; text-align: justify; }
  .md p:last-child { margin-bottom: 0; }
  .md a { color:#780000; text-decoration: none; word-break: break-all; }
  .md a:hover { text-decoration: underline; }
  .md strong { font-weight: 600; }
  .md em { font-style: italic; }

  table { width:100%; border-collapse:collapse; margin-top:8pt; }
  th, td { border-bottom:1px solid #e5e5e5; padding:6pt 8pt; text-align:left; font-size:9pt;
           vertical-align: top; line-height: 1.4; }
  th { background:#fafafa; font-weight:600; letter-spacing:0.02em;
       text-transform:uppercase; font-size:8.5pt; color:#6b6b6b; }
  .meta { color:#6b6b6b; font-size:9pt; }
  .meta-row { color:#6b6b6b; font-size:9pt; margin-top:2pt; }

  .pill-row { margin-top:8pt; }
  .pill { display:inline-block; padding:3pt 9pt; border-radius:11pt; font-size:8.5pt;
          background:#fafafa; border:1px solid #e5e5e5; margin-right:4pt; }
  .pill.maroon { background:#780000; color:#fff; border-color:#780000; font-weight:600; }
  .pill.green { background:#10b981; color:#fff; border-color:#10b981; font-weight:600; }
  .pill.red { background:#dc2626; color:#fff; border-color:#dc2626; font-weight:600; }
  .pill.amber { background:#f59e0b; color:#fff; border-color:#f59e0b; font-weight:600; }

  /* Severity-coloured rows */
  td.sev-CRITICAL { color:#dc2626; font-weight:600; }
  td.sev-HIGH { color:#dc2626; font-weight:500; }
  td.sev-MEDIUM { color:#f59e0b; font-weight:500; }
  td.sev-LOW { color:#10b981; font-weight:500; }

  ul, ol { margin-top:4pt; padding-left:18pt; }
  li { margin-bottom: 3pt; }
  .footer { color:#6b6b6b; font-size:8pt; margin-top:20pt; border-top:1px solid #e5e5e5;
            padding-top:6pt; }
  .version-tag { color:#6b6b6b; font-size:8pt; margin-top:4pt; }

  /* Verdict callout box */
  .verdict-box { border-left: 4px solid #780000; background:#fafafa;
                 padding: 9pt 12pt; margin-top: 8pt; }
  .verdict-box .label { color:#6b6b6b; font-size:8pt; text-transform:uppercase;
                        letter-spacing:0.08em; font-weight:600; }
</style></head>
<body>
<div class='brand-band'>
  {% if logo_data %}<img class='brand-logo' src='{{ logo_data }}' alt='Stone Gate'>{% endif %}
  <span class='label'>Investment Committee Memorandum</span>
</div>

<h1>{{ title }}</h1>
<p class='meta'>{{ sponsor }} &nbsp;·&nbsp; {{ city }}, {{ country }} &nbsp;·&nbsp;
   {{ property_type }} &nbsp;·&nbsp; {{ generated_at }}</p>
<p class='version-tag'>Analysis version v{{ analysis_version }}</p>

<p class='pill-row'>
  {% set verdict_class = 'maroon' %}
  {% if recommendation == 'PROCEED' %}{% set verdict_class = 'green' %}
  {% elif recommendation == 'REJECT' %}{% set verdict_class = 'red' %}
  {% elif recommendation == 'PROCEED_WITH_CONDITIONS' %}{% set verdict_class = 'amber' %}
  {% endif %}
  <span class='pill {{ verdict_class }}'>{{ recommendation }}</span>
  <span class='pill'>Opportunity {{ opportunity_score }} / 10</span>
  <span class='pill'>Risk {{ risk_score }} / 10</span>
  <span class='pill'>Confidence {{ confidence_score }} / 10</span>
</p>

<h2>Executive Summary</h2>
<div class='md'>{{ executive_summary_html }}</div>

<h2>Investment Thesis</h2>
<div class='md'>{{ thesis_html }}</div>

<h2>Deal Overview</h2>
<table>
<tr><th>Asking Equity</th><td>{{ asking_equity }}</td></tr>
<tr><th>Total Capitalization</th><td>{{ total_capitalization }}</td></tr>
<tr><th>Target IRR</th><td>{{ target_irr }}%</td></tr>
<tr><th>Target MOIC</th><td>{{ target_moic }}x</td></tr>
<tr><th>Hold Period</th><td>{{ hold_period_years }} yrs</td></tr>
</table>

{% if financial_analysis and (financial_analysis.verdict or financial_analysis.analysis or financial_analysis.headline_metrics or financial_analysis.sensitivity) %}
<h2>Financial Analysis</h2>

{% if financial_analysis.verdict or financial_analysis.verdict_rationale %}
<div class='verdict-box'>
  <div class='label'>Financial verdict</div>
  <p style='margin:4pt 0 0 0;'>
    {% set fv_class = 'maroon' %}
    {% if financial_analysis.verdict == 'STRONG' %}{% set fv_class = 'green' %}
    {% elif financial_analysis.verdict == 'MARGINAL' %}{% set fv_class = 'amber' %}
    {% elif financial_analysis.verdict == 'WEAK' %}{% set fv_class = 'red' %}
    {% endif %}
    <span class='pill {{ fv_class }}'>{{ financial_analysis.verdict or '—' }}</span>
  </p>
  {% if financial_analysis.verdict_rationale %}
  <div class='md' style='margin-top:6pt;'>{{ financial_verdict_rationale_html }}</div>
  {% endif %}
</div>
{% endif %}

{% if financial_analysis.headline_metrics %}
<h3>Headline metrics</h3>
<table>
<tr><th style='width:30%;'>Metric</th><th style='width:18%;'>Value</th><th>Interpretation</th></tr>
{% for m in financial_analysis.headline_metrics %}
<tr>
  <td><strong>{{ m.name }}</strong></td>
  <td class='tabular'><strong>{{ m.value }}</strong></td>
  <td>{{ m.interpretation }}</td>
</tr>
{% endfor %}
</table>
{% endif %}

{% if financial_analysis_narrative_html %}
<h3>Analysis</h3>
<div class='md'>{{ financial_analysis_narrative_html }}</div>
{% endif %}

{% if financial_analysis.sensitivity %}
<h3>Sensitivity — what breaks the deal</h3>
<table>
<tr><th>Variable</th><th>Base case</th><th>Breakpoint</th><th>Impact</th></tr>
{% for s in financial_analysis.sensitivity %}
<tr>
  <td><strong>{{ s.variable }}</strong></td>
  <td class='tabular'>{{ s.base_case or '—' }}</td>
  <td class='tabular' style='color:#dc2626;'>{{ s.breakpoint or '—' }}</td>
  <td>{{ s.impact or '—' }}</td>
</tr>
{% endfor %}
</table>
{% endif %}

{% if financial_analysis.weak_assumption_callouts %}
<h3>Weak-assumption callouts</h3>
<table>
<tr><th style='width:25%;'>Assumption</th><th style='width:15%;'>Sponsor</th><th style='width:12%;'>Severity</th><th>Rationale</th></tr>
{% for c in financial_analysis.weak_assumption_callouts %}
<tr>
  <td><strong>{{ c.name }}</strong></td>
  <td class='tabular'>{{ c.sponsor_value or '—' }}</td>
  <td class='sev-{{ c.severity or "MEDIUM" }}'>{{ c.severity or 'MEDIUM' }}</td>
  <td>{{ c.rationale }}</td>
</tr>
{% endfor %}
</table>
{% endif %}
{% endif %}

<h2>Market Analysis</h2>
<div class='md'>{{ market_analysis_html }}</div>

<h2>Risk Analysis</h2>
<table><tr><th>Category</th><th>Severity</th><th>Risk</th><th>Mitigation</th></tr>
{% for r in risks %}
<tr>
  <td>{{ r.category }}</td>
  <td class='sev-{{ r.severity }}'>{{ r.severity }}</td>
  <td><strong>{{ r.title }}</strong><div class='meta-row'>{{ r.description }}</div></td>
  <td>{{ r.mitigation or '—' }}</td>
</tr>
{% endfor %}
</table>

<h2>Bull / Base / Bear</h2>
<h3>Bull case</h3><ul>{% for x in bull_case %}<li>{{ x }}</li>{% endfor %}</ul>
<h3>Base case</h3><ul>{% for x in base_case %}<li>{{ x }}</li>{% endfor %}</ul>
<h3>Bear case</h3><ul>{% for x in bear_case %}<li>{{ x }}</li>{% endfor %}</ul>

<h2>Recommendation</h2>
<div class='verdict-box'>
  <div class='label'>AI Verdict</div>
  <p style='margin:4pt 0 0 0;'><strong>{{ ai_verdict }}</strong></p>
  {% if ai_verdict_rationale_html %}<div class='md' style='margin-top:6pt;'>{{ ai_verdict_rationale_html }}</div>{% endif %}
</div>

{% if ai_top_reasons_for %}
<h3>Reasons in favor</h3>
<ul>{% for r in ai_top_reasons_for %}<li>{{ r }}</li>{% endfor %}</ul>
{% endif %}

{% if ai_top_reasons_against %}
<h3>Reasons against</h3>
<ul>{% for r in ai_top_reasons_against %}<li>{{ r }}</li>{% endfor %}</ul>
{% endif %}

{% if ai_critical_questions %}
<h3>Decision-critical questions</h3>
<ol>{% for q in ai_critical_questions %}<li>{{ q }}</li>{% endfor %}</ol>
{% endif %}

{% if ai_next_steps %}
<h3>Next steps</h3>
<ul>{% for s in ai_next_steps %}<li>{{ s }}</li>{% endfor %}</ul>
{% endif %}

{% if ai_watchpoints %}
<h3>Watchpoints (if approved)</h3>
<ul>{% for w in ai_watchpoints %}<li>{{ w }}</li>{% endfor %}</ul>
{% endif %}

<div class='footer'>
Prepared by Stone Gate Investment Intelligence System.
This document is confidential and contains AI-generated analysis subject to human review.
</div>
</body></html>
"""

EXEC_SUMMARY_TEMPLATE = """\
<!doctype html><html><head><meta charset='utf-8'><style>
  @page { size: A4; margin: 22mm 18mm; }
  body { font-family:'Inter','Helvetica Neue',sans-serif; color:#0a0a0a;
         font-size:10pt; line-height:1.55; }
  .brand-band { background:#780000; color:#fff; padding:8pt 14pt;
                display:flex; align-items:center; justify-content:space-between;
                margin:-8pt -8pt 14pt -8pt; }
  .brand-logo { height: 20pt; width:auto; mix-blend-mode: screen; }
  .brand-band .label { font-weight:600; text-transform:uppercase; font-size:8.5pt;
                       letter-spacing:.08em; color:#fff; }
  h1 { font-size:18pt; margin-top:6pt; margin-bottom:4pt; font-weight:700; }
  h2 { font-size:11.5pt; color:#780000; margin-top:16pt; margin-bottom:6pt;
       font-weight:600; }
  .meta { color:#6b6b6b; font-size:9pt; }
  .pill-row { margin-top:8pt; }
  .pill { display:inline-block; padding:3pt 9pt; border-radius:11pt; font-size:8.5pt;
          background:#fafafa; border:1px solid #e5e5e5; margin-right:4pt; }
  .pill.maroon { background:#780000; color:#fff; border-color:#780000; font-weight:600; }
  .pill.green { background:#10b981; color:#fff; border-color:#10b981; font-weight:600; }
  .pill.red { background:#dc2626; color:#fff; border-color:#dc2626; font-weight:600; }
  .pill.amber { background:#f59e0b; color:#fff; border-color:#f59e0b; font-weight:600; }
  .md p { margin: 0 0 8pt 0; text-align: justify; }
  .md p:last-child { margin-bottom: 0; }
  .md a { color:#780000; text-decoration: none; word-break: break-all; }
  .md strong { font-weight: 600; }
  .md em { font-style: italic; }
  ol li { margin-bottom: 5pt; }
  .footer { color:#6b6b6b; font-size:8pt; margin-top:20pt; border-top:1px solid #e5e5e5;
            padding-top:6pt; }
</style></head><body>
<div class='brand-band'>
  {% if logo_data %}<img class='brand-logo' src='{{ logo_data }}' alt='Stone Gate'>{% endif %}
  <span class='label'>Executive Summary</span>
</div>
<h1>{{ title }}</h1>
<p class='meta'>{{ sponsor }} · {{ city }}, {{ country }} · {{ property_type }}</p>
<p class='pill-row'>
  {% set verdict_class = 'maroon' %}
  {% if recommendation == 'PROCEED' %}{% set verdict_class = 'green' %}
  {% elif recommendation == 'REJECT' %}{% set verdict_class = 'red' %}
  {% elif recommendation == 'PROCEED_WITH_CONDITIONS' %}{% set verdict_class = 'amber' %}
  {% endif %}
  <span class='pill {{ verdict_class }}'>{{ recommendation }}</span>
</p>
<h2>Snapshot</h2><div class='md'>{{ executive_summary_html }}</div>
<h2>Why now</h2><div class='md'>{{ thesis_html }}</div>

{% if financial_analysis and (financial_analysis.verdict or financial_analysis.headline_metrics or financial_analysis.sensitivity) %}
<h2>Financial snapshot</h2>
{% if financial_analysis.verdict or financial_analysis.verdict_rationale %}
<p>
  {% set fv_class = 'maroon' %}
  {% if financial_analysis.verdict == 'STRONG' %}{% set fv_class = 'green' %}
  {% elif financial_analysis.verdict == 'MARGINAL' %}{% set fv_class = 'amber' %}
  {% elif financial_analysis.verdict == 'WEAK' %}{% set fv_class = 'red' %}
  {% endif %}
  <span class='pill {{ fv_class }}'>{{ financial_analysis.verdict or '—' }}</span>
</p>
{% if financial_analysis.verdict_rationale %}
<div class='md'>{{ financial_verdict_rationale_html }}</div>
{% endif %}
{% endif %}
{% if financial_analysis.headline_metrics %}
<ul style='margin-top:6pt;'>
{% for m in financial_analysis.headline_metrics[:4] %}
<li><strong>{{ m.name }}</strong> — {{ m.value }}. {{ m.interpretation }}</li>
{% endfor %}
</ul>
{% endif %}
{% if financial_analysis.sensitivity %}
<p class='meta' style='margin-top:6pt;'><strong>Key sensitivities:</strong></p>
<ul>
{% for s in financial_analysis.sensitivity[:3] %}
<li><strong>{{ s.variable }}</strong> ({{ s.base_case or '—' }} → {{ s.breakpoint or '—' }}) — {{ s.impact or '' }}</li>
{% endfor %}
</ul>
{% endif %}
{% endif %}

<h2>Top three risks</h2>
<ol>{% for r in risks[:3] %}<li><strong>{{ r.title }}</strong> — {{ r.description }}</li>{% endfor %}</ol>
<h2>Recommendation</h2><div class='md'>{{ recommendation_rationale_html }}</div>
<div class='footer'>Stone Gate · Confidential · AI-generated, human-reviewed</div>
</body></html>
"""


def _render_html(template: str, payload: dict, logo_path: Path = LOGO_PATH) -> str:
    env = Environment(loader=BaseLoader(), autoescape=True, trim_blocks=True, lstrip_blocks=True)
    return env.from_string(template).render(**payload, logo_data=_logo_data_uri(logo_path))


def _enrich_payload_with_html(payload: dict) -> dict:
    """Convert markdown-bearing prose fields to safe HTML so the templates
    can render them with proper headings, links, emphasis."""
    enriched = dict(payload)
    enriched["executive_summary_html"] = md_to_html(payload.get("executive_summary"))
    enriched["thesis_html"] = md_to_html(payload.get("thesis"))
    enriched["market_analysis_html"] = md_to_html(payload.get("market_analysis"))
    enriched["ai_verdict_rationale_html"] = md_to_html(payload.get("ai_verdict_rationale"))
    enriched["recommendation_rationale_html"] = md_to_html(
        payload.get("recommendation_rationale")
    )
    # Financial analyst's narrative + verdict rationale carry markdown too.
    fa = payload.get("financial_analysis") or {}
    enriched["financial_analysis_narrative_html"] = md_to_html(fa.get("analysis"))
    enriched["financial_verdict_rationale_html"] = md_to_html(fa.get("verdict_rationale"))
    return enriched


def render_pdf(payload: dict, kind: str, out_dir: str) -> str:
    """Render PDF with WeasyPrint. Returns the output path."""
    from weasyprint import HTML

    if kind == "IC_MEMO_LONG":
        template = IC_MEMO_TEMPLATE
        logo_path = LOGO_PATH
    else:
        template = EXEC_SUMMARY_TEMPLATE
        # Executive Summary uses the white-on-transparent logo so it sits
        # cleanly on the maroon header band (no white rectangle behind it).
        logo_path = LOGO_WHITE_PATH if LOGO_WHITE_PATH.exists() else LOGO_PATH
    html = _render_html(template, _enrich_payload_with_html(payload), logo_path)
    Path(out_dir).mkdir(parents=True, exist_ok=True)
    out = os.path.join(out_dir, f"{uuid.uuid4()}.pdf")
    HTML(string=html).write_pdf(out)
    return out


def _add_content_chrome(slide, prs, slide_num: int) -> None:
    """Content-slide branding: left vertical maroon bar + Stone Gate logo in
    bottom-right corner + page number. Title is rendered separately by the
    caller (with its own maroon underline). Used on all slides AFTER the
    title slide."""
    # Left vertical maroon bar (full height, narrow)
    bar_w = Inches(0.3)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), bar_w, prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY_RGB
    bar.line.fill.background()

    # Stone Gate logo bottom-right (maroon-on-white version reads on white slide)
    if LOGO_PATH.exists():
        logo_h = Inches(0.55)
        # logo width is roughly 3.5x its height for this script
        logo_w_estimate = Inches(1.9)
        slide.shapes.add_picture(
            str(LOGO_PATH),
            prs.slide_width - logo_w_estimate - Inches(0.4),
            prs.slide_height - logo_h - Inches(0.25),
            height=logo_h,
        )

    # Page number — small, muted grey, bottom-right
    pn = slide.shapes.add_textbox(
        prs.slide_width - Inches(0.45),
        prs.slide_height - Inches(0.45),
        Inches(0.3),
        Inches(0.25),
    ).text_frame
    pn.text = str(slide_num)
    pn.paragraphs[0].font.size = Pt(9)
    pn.paragraphs[0].font.color.rgb = RGBColor(0x6B, 0x6B, 0x6B)
    pn.paragraphs[0].alignment = 2  # right


def _add_title_with_rule(slide, title: str, subtitle: str | None = None) -> float:
    """Add the maroon title and a maroon underline to a content slide. Returns
    the Y offset (in inches) where body content should start."""
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.7)).text_frame
    tb.text = title
    tb.paragraphs[0].font.size = Pt(28)
    tb.paragraphs[0].font.bold = True
    tb.paragraphs[0].font.color.rgb = PRIMARY_RGB

    # Maroon underline rule under the title
    rule = slide.shapes.add_shape(
        1,
        Inches(0.7),
        Inches(1.18),
        Inches(12.0),
        Inches(0.035),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = PRIMARY_RGB
    rule.line.fill.background()

    y = 1.35
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(12), Inches(0.35)).text_frame
        sb.text = subtitle
        sb.paragraphs[0].font.size = Pt(11)
        sb.paragraphs[0].font.color.rgb = RGBColor(0x6B, 0x6B, 0x6B)
        y += 0.4
    return y + 0.2


def _add_brand_band(slide, prs, label_text: str | None = None) -> None:
    """Add a maroon top band with the Stone Gate logo on the left and an
    optional uppercase label on the right."""
    band_h = Inches(0.55)
    band = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, band_h)
    band.fill.solid()
    band.fill.fore_color.rgb = PRIMARY_RGB
    band.line.fill.background()

    if LOGO_PATH.exists():
        # Place logo at left. python-pptx doesn't have mix-blend-mode, so we accept that
        # the logo's black background fills part of the band — but it sits on the
        # maroon band so the white script is what reads. The black bg is a small visual
        # rectangle on the maroon — acceptable until we ship a transparent-bg logo.
        slide.shapes.add_picture(
            str(LOGO_PATH),
            Inches(0.25),
            Inches(0.07),
            height=Inches(0.42),
        )

    if label_text:
        right_w = Inches(5)
        right_x = prs.slide_width - right_w - Inches(0.3)
        tb = slide.shapes.add_textbox(right_x, Inches(0.1), right_w, Inches(0.4)).text_frame
        tb.text = label_text
        p = tb.paragraphs[0]
        p.alignment = 2  # right
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE_RGB


# ─── PPTX helpers: concise bullets from prose ───────────────────────

def _strip_md(text: str) -> str:
    """Strip markdown markers for clean slide rendering."""
    if not text:
        return ""
    t = re.sub(r"\*\*([^*]+?)\*\*", r"\1", text)
    t = re.sub(r"\*([^*\s][^*]*?)\*", r"\1", t)
    t = re.sub(r"\[([^\]]+)\]\(https?://[^)]+\)", r"\1", t)  # link label only
    t = re.sub(r"https?://\S+", "", t)  # strip bare URLs
    t = re.sub(r"\s+", " ", t).strip()
    return t


def _prose_to_bullets(
    text: str | None, max_bullets: int = 5, max_chars: int = 180
) -> list[str]:
    """Extract concise bullets from prose. Takes the first sentence of each
    paragraph; truncates to max_chars; caps at max_bullets total."""
    if not text:
        return []
    # Skip AI preamble + drop any '# Title' line
    cleaned = _strip_ai_preamble(text)
    cleaned = re.sub(r"^#\s+[^\n]*\n+", "", cleaned)
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", cleaned) if p.strip()]
    bullets: list[str] = []
    for p in paragraphs:
        # Skip headings — already represented as section structure
        if p.startswith("#"):
            continue
        first = re.split(r"(?<=[.!?])\s+", p)[0]
        first = _strip_md(first)
        if not first:
            continue
        if len(first) > max_chars:
            first = first[: max_chars - 1].rstrip() + "…"
        bullets.append(first)
        if len(bullets) >= max_bullets:
            break
    return bullets


def _truncate(text: str | None, max_chars: int) -> str:
    s = _strip_md(text or "")
    if len(s) <= max_chars:
        return s
    return s[: max_chars - 1].rstrip() + "…"


def render_pptx(payload: dict, out_dir: str) -> str:
    Path(out_dir).mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # ── Title slide ─────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    _add_brand_band(slide, prs, label_text=None)
    if LOGO_PATH.exists():
        # Big centered logo on the title slide too
        slide.shapes.add_picture(
            str(LOGO_PATH),
            Inches(4.6), Inches(1.4),
            height=Inches(1.8),
        )
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(3.6), Inches(12), Inches(2)).text_frame
    tb.text = payload.get("title", "Stone Gate IC Presentation")
    tb.paragraphs[0].font.size = Pt(36)
    tb.paragraphs[0].font.bold = True
    tb.paragraphs[0].alignment = 1  # center
    sub = tb.add_paragraph()
    sub.text = (
        f"{payload.get('sponsor','')} · {payload.get('city','')}, {payload.get('country','')}"
    )
    sub.alignment = 1
    sub.font.size = Pt(16)
    sub.font.color.rgb = RGBColor(0x6B, 0x6B, 0x6B)

    rec = slide.shapes.add_textbox(Inches(0.6), Inches(5.6), Inches(12), Inches(0.5)).text_frame
    rec.text = f"Recommendation: {payload.get('recommendation', 'NEED MORE INFO')}"
    rec.paragraphs[0].alignment = 1
    rec.paragraphs[0].font.size = Pt(13)
    rec.paragraphs[0].font.color.rgb = PRIMARY_RGB
    rec.paragraphs[0].font.bold = True

    # ── Section slides ──────────────────────────────────────
    # Slide counter is shared by all add_*_slide builders below so the
    # bottom-right page number stays continuous.
    counter = {"n": 1}  # title slide is page 1 (we don't number it)

    def _next_num() -> int:
        counter["n"] += 1
        return counter["n"]

    # Content slides start at x=0.7" (just inside the left maroon bar) and
    # extend to ~12.5" wide. Body area below the title rule sits at ~1.7".
    BODY_LEFT = 0.7
    BODY_WIDTH = 12.0

    def add_content_slide(title: str, bullets: list[str], subtitle: str | None = None):
        slide = prs.slides.add_slide(blank)
        _add_content_chrome(slide, prs, _next_num())
        body_top = _add_title_with_rule(slide, title, subtitle)
        body = slide.shapes.add_textbox(
            Inches(BODY_LEFT), Inches(body_top), Inches(BODY_WIDTH),
            Inches(7.5 - body_top - 0.9),  # leave room for logo + page num
        ).text_frame
        body.word_wrap = True
        if not bullets:
            p = body.paragraphs[0]
            p.text = "—"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
            return
        for i, b in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = f"•  {b}"
            p.font.size = Pt(14)
            p.space_after = Pt(7)

    def add_two_col_slide(
        title: str, left_label: str, left_bullets: list[str],
        right_label: str, right_bullets: list[str],
    ):
        slide = prs.slides.add_slide(blank)
        _add_content_chrome(slide, prs, _next_num())
        body_top = _add_title_with_rule(slide, title)
        col_width = Inches(BODY_WIDTH / 2 - 0.1)
        col_height = Inches(7.5 - body_top - 1.0)

        for col_idx, (label, bullets) in enumerate(
            [(left_label, left_bullets), (right_label, right_bullets)]
        ):
            x = Inches(BODY_LEFT + col_idx * (BODY_WIDTH / 2 + 0.05))
            lab = slide.shapes.add_textbox(x, Inches(body_top), col_width, Inches(0.35)).text_frame
            lab.text = label
            lab.paragraphs[0].font.size = Pt(11)
            lab.paragraphs[0].font.bold = True
            lab.paragraphs[0].font.color.rgb = PRIMARY_RGB
            body = slide.shapes.add_textbox(
                x, Inches(body_top + 0.4), col_width, col_height
            ).text_frame
            body.word_wrap = True
            if not bullets:
                p = body.paragraphs[0]
                p.text = "—"
                p.font.size = Pt(13)
                p.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
                continue
            for i, b in enumerate(bullets):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                p.text = f"•  {b}"
                p.font.size = Pt(13)
                p.space_after = Pt(6)

    def add_three_col_slide(
        title: str, cols: list[tuple[str, list[str]]],
    ):
        slide = prs.slides.add_slide(blank)
        _add_content_chrome(slide, prs, _next_num())
        body_top = _add_title_with_rule(slide, title)
        col_width = Inches(BODY_WIDTH / 3 - 0.1)
        col_height = Inches(7.5 - body_top - 1.0)

        col_colors = [
            RGBColor(0x10, 0xB9, 0x81),  # green for bull
            RGBColor(0x6B, 0x6B, 0x6B),  # neutral grey for base
            RGBColor(0xDC, 0x26, 0x26),  # red for bear
        ]
        for col_idx, (label, bullets) in enumerate(cols):
            x = Inches(BODY_LEFT + col_idx * (BODY_WIDTH / 3 + 0.05))
            lab = slide.shapes.add_textbox(x, Inches(body_top), col_width, Inches(0.35)).text_frame
            lab.text = label
            lab.paragraphs[0].font.size = Pt(12)
            lab.paragraphs[0].font.bold = True
            lab.paragraphs[0].font.color.rgb = (
                col_colors[col_idx] if col_idx < len(col_colors) else PRIMARY_RGB
            )
            body = slide.shapes.add_textbox(
                x, Inches(body_top + 0.4), col_width, col_height
            ).text_frame
            body.word_wrap = True
            if not bullets:
                p = body.paragraphs[0]
                p.text = "—"
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
                continue
            for i, b in enumerate(bullets):
                p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                p.text = f"•  {b}"
                p.font.size = Pt(12)
                p.space_after = Pt(6)

    # ── 1. Executive snapshot ──────────────────────────────────
    add_content_slide(
        "Executive snapshot",
        _prose_to_bullets(payload.get("executive_summary"), max_bullets=5),
        subtitle=f"{payload.get('sponsor','')} · {payload.get('city','')}, {payload.get('country','')} · {payload.get('property_type','')}",
    )

    # ── 2. Deal economics ──────────────────────────────────────
    add_content_slide(
        "Deal economics",
        [
            f"Asking equity:  {payload.get('asking_equity', '—')}",
            f"Total capitalization:  {payload.get('total_capitalization', '—')}",
            f"Target IRR:  {payload.get('target_irr', '—')}%",
            f"Target MOIC:  {payload.get('target_moic', '—')}x",
            f"Hold period:  {payload.get('hold_period_years', '—')} years",
            f"AI scores:  Opportunity {payload.get('opportunity_score','—')} / Risk {payload.get('risk_score','—')} / Confidence {payload.get('confidence_score','—')}",
        ],
    )

    # ── 3. Investment thesis ───────────────────────────────────
    add_content_slide(
        "Investment thesis",
        _prose_to_bullets(payload.get("thesis"), max_bullets=5, max_chars=200),
    )

    # ── 4. Market research ─────────────────────────────────────
    market_text = payload.get("market_analysis") or ""
    market_verdict = ""
    m_match = re.search(r"\b(TAILWIND|NEUTRAL|HEADWIND)\b", market_text)
    if m_match:
        market_verdict = m_match.group(1)
    market_bullets = _prose_to_bullets(market_text, max_bullets=5, max_chars=200)
    add_content_slide(
        "Market research",
        market_bullets,
        subtitle=f"Verdict: {market_verdict}" if market_verdict else None,
    )

    # ── 5. Financial verdict + headline metrics ────────────────
    fa = payload.get("financial_analysis") or {}
    fa_bullets: list[str] = []
    if fa.get("verdict_rationale"):
        fa_bullets.append(_truncate(fa["verdict_rationale"], 240))
    for m in (fa.get("headline_metrics") or [])[:6]:
        line = f"{m.get('name','')}: {m.get('value','')}"
        if m.get("interpretation"):
            line += f" — {_truncate(m.get('interpretation'), 130)}"
        fa_bullets.append(line)
    add_content_slide(
        "Financial verdict",
        fa_bullets,
        subtitle=f"Verdict: {fa.get('verdict', '—')}",
    )

    # ── 6. Sensitivities ───────────────────────────────────────
    sens_bullets = []
    for s in (fa.get("sensitivity") or [])[:5]:
        base = s.get("base_case") or "—"
        breakpoint = s.get("breakpoint") or "—"
        impact = _truncate(s.get("impact"), 130)
        sens_bullets.append(f"{s.get('variable','')}  ({base} → {breakpoint})  — {impact}")
    if sens_bullets:
        add_content_slide(
            "Sensitivities — what breaks the deal",
            sens_bullets,
        )

    # ── 7. Weak assumption callouts ────────────────────────────
    weak_bullets = []
    for c in (fa.get("weak_assumption_callouts") or [])[:5]:
        sev = c.get("severity") or "MEDIUM"
        line = f"[{sev}] {c.get('name','')}"
        if c.get("sponsor_value"):
            line += f"  (sponsor: {c['sponsor_value']})"
        if c.get("rationale"):
            line += f" — {_truncate(c.get('rationale'), 140)}"
        weak_bullets.append(line)
    if weak_bullets:
        add_content_slide("Weak-assumption callouts", weak_bullets)

    # ── 8. Top risks ───────────────────────────────────────────
    risks_payload = payload.get("risks") or []
    add_content_slide(
        "Top risks",
        [
            f"[{r.get('severity','MEDIUM')}] {r.get('title','')} — {_truncate(r.get('description'), 130)}"
            for r in risks_payload[:6]
        ],
    )

    # ── 9. Bull / Base / Bear (three columns) ──────────────────
    bull = [_truncate(b, 180) for b in (payload.get("bull_case") or [])[:5]]
    base = [_truncate(b, 180) for b in (payload.get("base_case") or [])[:5]]
    bear = [_truncate(b, 180) for b in (payload.get("bear_case") or [])[:5]]
    add_three_col_slide(
        "Bull · Base · Bear",
        [("Bull case", bull), ("Base case", base), ("Bear case", bear)],
    )

    # ── 10. AI verdict + reasons for / against ─────────────────
    for_bullets = [_truncate(r, 200) for r in (payload.get("ai_top_reasons_for") or [])[:5]]
    against_bullets = [_truncate(r, 200) for r in (payload.get("ai_top_reasons_against") or [])[:5]]
    add_two_col_slide(
        f"AI verdict: {payload.get('ai_verdict','NEED_MORE_INFO')}",
        "Reasons in favor", for_bullets,
        "Reasons against", against_bullets,
    )

    # ── 11. Critical questions for IC ──────────────────────────
    crit = [_truncate(q, 200) for q in (payload.get("ai_critical_questions") or [])[:6]]
    if crit:
        add_content_slide("Decision-critical questions", crit)

    # ── 12. Next steps + watchpoints (two columns) ─────────────
    nxt = [_truncate(s, 180) for s in (payload.get("ai_next_steps") or [])[:6]]
    watch = [_truncate(w, 180) for w in (payload.get("ai_watchpoints") or [])[:6]]
    add_two_col_slide(
        "Next steps · Watchpoints",
        "Next steps (pre-IC)", nxt,
        "Watchpoints (if approved)", watch,
    )

    # ── 13. Recommendation closer ──────────────────────────────
    rec_bullets = []
    rationale = payload.get("recommendation_rationale") or payload.get("ai_verdict_rationale")
    if rationale:
        rec_bullets.extend(_prose_to_bullets(rationale, max_bullets=4, max_chars=220))
    add_content_slide(
        f"Recommendation: {payload.get('recommendation', 'NEED MORE INFO')}",
        rec_bullets,
    )

    out = os.path.join(out_dir, f"{uuid.uuid4()}.pptx")
    prs.save(out)
    return out
