"""One-shot generator for Stone Gate platform documentation.

Produces two files using the same reporting infrastructure the IC memos use:
  - stone-gate-architecture.pdf  — technical architecture deep-dive (WeasyPrint)
  - stone-gate-overview.pptx     — non-technical platform overview (python-pptx)

Run from inside the ai-service container:
    docker exec stonegate-ai python -m app.scripts.generate_platform_docs

Output is written to /app/data/blobs/platform-docs/ which is bind-mounted
to <project>/data/blobs/platform-docs/ on the host.
"""
from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from weasyprint import HTML
from jinja2 import Environment, BaseLoader

from app.reporting import (
    PRIMARY_RGB,
    LOGO_PATH,
    _logo_data_uri,
    _add_content_chrome,
    _add_title_with_rule,
)


# ─── Technical Architecture PDF ──────────────────────────────────────

ARCH_HTML = r"""<!doctype html>
<html><head><meta charset='utf-8'><style>
  @page { size: A4; margin: 32mm 0 22mm 0;
          @top-left-corner { content: ""; }
          @top-left { content: element(header); width: 210mm; padding: 0; margin: 0; }
          @top-right-corner { content: ""; }
          @bottom-right { content: "Page " counter(page) " of " counter(pages);
                          color:#6b6b6b; font-size:8pt; margin-right: 18mm; } }
  body { font-family: 'Inter','Helvetica Neue',sans-serif; color:#0a0a0a;
         font-size: 10pt; line-height: 1.55; padding: 0 18mm; }
  h1, h2, h3, h4 { color:#0a0a0a; }
  .brand-band { position: running(header); background:#ffffff;
                border-bottom: 1.5pt solid #A88B47;
                padding: 10pt 18mm;
                display: flex; align-items: center; gap: 14pt;
                justify-content: flex-start; }
  .brand-logo { height: 44pt; width: auto; }
  .brand-band .label { font-weight: 600; letter-spacing: 0.08em; text-transform: uppercase;
                       font-size: 9.5pt; color: #A88B47; }
  h1 { font-size: 26pt; margin-top: 4pt; margin-bottom: 4pt; font-weight: 700; }
  h2 { font-size: 14pt; border-bottom: 1.5px solid #A88B47; padding-bottom: 4pt;
       margin-top: 20pt; margin-bottom: 8pt; font-weight: 600; }
  h3 { font-size: 11.5pt; color: #A88B47; margin-top: 14pt; margin-bottom: 6pt; }
  p { text-align: justify; margin: 0 0 8pt 0; }
  table { width: 100%; border-collapse: collapse; margin-top: 6pt; }
  th, td { border-bottom: 1px solid #e5e5e5; padding: 6pt 8pt; text-align: left;
           font-size: 9pt; vertical-align: top; line-height: 1.4; }
  th { background: #fafafa; font-weight: 600; letter-spacing: 0.02em;
       text-transform: uppercase; font-size: 8.5pt; color: #6b6b6b; }
  .meta { color: #6b6b6b; font-size: 9pt; }
  ul, ol { padding-left: 18pt; margin-top: 4pt; }
  li { margin-bottom: 3pt; }
  pre { background: #fafafa; border: 1px solid #e5e5e5; padding: 8pt 10pt;
        font-family: 'Consolas','Courier New',monospace; font-size: 8.5pt;
        line-height: 1.4; overflow-x: hidden; white-space: pre-wrap;
        margin-top: 6pt; margin-bottom: 8pt; }
  .pill { display: inline-block; padding: 2pt 8pt; border-radius: 10pt;
          font-size: 8.5pt; background: #fafafa; border: 1px solid #e5e5e5;
          margin-right: 3pt; }
  .pill.maroon { background: #A88B47; color: #fff; border-color: #A88B47; }
  .pill.green  { background: #10b981; color: #fff; border-color: #10b981; }
  .pill.amber  { background: #f59e0b; color: #fff; border-color: #f59e0b; }
  .pill.red    { background: #dc2626; color: #fff; border-color: #dc2626; }
  .toc { margin: 14pt 0; }
  .toc li { margin-bottom: 4pt; }
  .cover { page-break-after: always; padding-top: 80pt; }
  .cover h1 { font-size: 40pt; line-height: 1.1; color: #A88B47; }
  .cover .sub { font-size: 14pt; color: #6b6b6b; margin-top: 14pt; }
  .cover .meta { margin-top: 32pt; font-size: 10pt; color: #6b6b6b; }
</style></head><body>

<div class='brand-band'>
  {% if logo_data %}<img class='brand-logo' src='{{ logo_data }}' alt='Stone Gate'>{% endif %}
  <span class='label'>Platform Architecture</span>
</div>

<!-- COVER -->
<div class='cover'>
  <h1>Stone Gate<br>Platform Architecture</h1>
  <div class='sub'>An AI-native investment diligence operating system</div>
  <div class='meta'>
    Technical reference · Version 1.0<br>
    {{ generated_at }}
  </div>
</div>

<!-- TOC -->
<h2>Contents</h2>
<ol class='toc'>
  <li>Executive summary</li>
  <li>System topology</li>
  <li>Frontend &amp; user interface</li>
  <li>API gateway</li>
  <li>AI service</li>
  <li>Document processor</li>
  <li>Data layer</li>
  <li>AI agent pipeline</li>
  <li>End-to-end flow: Analyze an opportunity</li>
  <li>Security model</li>
  <li>Backup &amp; recovery</li>
  <li>Deployment model</li>
</ol>

<!-- 1. EXECUTIVE SUMMARY -->
<h2>1. Executive summary</h2>
<p>Stone Gate is a local-first, AI-native platform for institutional real estate diligence.
It turns a fragmented data room (PDFs, Excel models, decks, emails) into a defensible
investment recommendation through a structured workflow: <strong>Context → Initial Screening →
Understanding → CIO Review → Recommendation → Monitoring</strong>.</p>
<p>The architecture is a monorepo of six runtime components:</p>
<ul>
  <li>A Next.js web application (the analyst workspace)</li>
  <li>A Fastify Node API (the trust boundary and audit log)</li>
  <li>A FastAPI AI service (agent orchestration, RAG, financial engine, scenario engine, reporting)</li>
  <li>A FastAPI document processor (PDF / Excel / Word / image / email ingestion + OCR)</li>
  <li>PostgreSQL with pgvector (relational + embedding store)</li>
  <li>ChromaDB and Redis (per-opportunity vector cache and task queue)</li>
</ul>
<p>The AI layer is a six-agent pipeline backed by Anthropic Claude with prompt
caching enabled. Four agents run in parallel; two run sequentially with cross-agent
context. The system produces a synthesis verdict (PROCEED / PROCEED_WITH_CONDITIONS /
REJECT / NEED_MORE_INFO) plus three report formats (IC memorandum PDF, executive
summary PDF, IC presentation deck).</p>

<!-- 2. SYSTEM TOPOLOGY -->
<h2>2. System topology</h2>
<pre>
┌────────────────────────────────────────────────────────────────────────┐
│                          Stone Gate Platform                           │
│                                                                        │
│  ┌──────────────────┐    ┌──────────────────┐                          │
│  │  Next.js Web     │◄──►│  Node API (BFF)  │                          │
│  │  (App Router)    │    │  Fastify + JWT   │                          │
│  │  :3100           │    │  :4000           │                          │
│  └──────────────────┘    └────────┬─────────┘                          │
│                                   │                                    │
│              ┌────────────────────┼────────────────────┐               │
│              ▼                    ▼                    ▼               │
│      ┌──────────────┐     ┌──────────────┐     ┌──────────────┐        │
│      │ AI Service   │     │ Doc          │     │ PostgreSQL + │        │
│      │ FastAPI      │     │ Processor    │     │ pgvector     │        │
│      │ Anthropic    │     │ FastAPI      │     │ :5432        │        │
│      │ :8000        │     │ :8003        │     │              │        │
│      └──────┬───────┘     └──────┬───────┘     └──────────────┘        │
│             │                    │                                     │
│             ▼                    ▼                                     │
│      ┌──────────────┐     ┌──────────────┐                             │
│      │ ChromaDB     │     │ Redis        │                             │
│      │ :8002        │     │ :6379        │                             │
│      └──────────────┘     └──────────────┘                             │
│                                                                        │
│      ┌────────────────────────────────────────────────┐                │
│      │ Local encrypted blob store  (./data/blobs)     │                │
│      │ AES-256-GCM, derived from MASTER_ENCRYPTION_KEY│                │
│      └────────────────────────────────────────────────┘                │
└────────────────────────────────────────────────────────────────────────┘
</pre>
<p>Every cross-service call passes through the Node API. Service-to-service
authentication uses short-lived signed tokens. The Node API is the only
component the browser ever talks to directly.</p>

<!-- 3. FRONTEND -->
<h2>3. Frontend &amp; user interface</h2>
<p>The web application is built on Next.js 14 (App Router) with server components for
data fetching and client components for interactivity (chat streaming, charts, filters).
Tailwind CSS is configured with Stone Gate design tokens; shadcn/ui primitives are
extended into domain-specific components (<code>MetricCard</code>, <code>RiskHeatmap</code>,
<code>QuadrantPanel</code>, <code>CollapsibleSection</code>, <code>ScenarioPanel</code>).</p>
<h3>Routing</h3>
<p>The opportunity workspace is organized as tabbed sub-routes following the workflow:
<code>/opportunities/[id]/{context,initial-screening,understanding,cio-review,recommendation,monitoring}</code>.
A computed <code>workflowStage</code> field on each opportunity drives the stepper in the
header and the kanban on the pipeline page; the value is derived from the opportunity's
data (does it have a thesis? a decision? approved by IC?) rather than stored as a column.</p>
<h3>State</h3>
<p>TanStack Query manages the server cache for opportunities, documents, threads.
Per-route React state handles the rest (filters, collapsible open/closed). JWT
authentication tokens are kept in <code>localStorage</code> with a 12-hour TTL. The Node
API is reachable cross-origin via a <code>Bearer</code> header.</p>

<!-- 4. API GATEWAY -->
<h2>4. API gateway</h2>
<p>The Node API in <code>apps/api</code> is the trust boundary. It enforces JWT
authentication on all routes (except <code>/v1/auth</code>), records every state-changing
operation to the <code>AuditLog</code> table, and forwards work to the Python services.</p>
<table>
<tr><th>Module</th><th>Responsibility</th></tr>
<tr><td>routes/auth</td><td>Login, logout, /me, password rotation</td></tr>
<tr><td>routes/opportunities</td><td>CRUD, briefing, analyze trigger, financial-only re-run, decision</td></tr>
<tr><td>routes/documents</td><td>Upload via multipart, processing status, soft-delete</td></tr>
<tr><td>routes/chat</td><td>Threads, message persistence, streaming proxy</td></tr>
<tr><td>routes/scenarios</td><td>Scenario CRUD, run trigger</td></tr>
<tr><td>routes/reports</td><td>Generate (IC memo / exec summary / deck), authenticated download</td></tr>
<tr><td>routes/telemetry</td><td>Aggregated LLM call counts, token usage, cost, latency</td></tr>
<tr><td>middleware/audit</td><td>Decorates every request with <code>req.audit(...)</code></td></tr>
</table>
<p>File uploads are streamed via <code>@fastify/multipart</code> with a 200 MB per-file
ceiling. Uploaded blobs are AES-256-GCM encrypted at rest using a key derived from
<code>MASTER_ENCRYPTION_KEY</code>; only relative paths are stored in the database so
the ai-service and Node API can each resolve them against their own
<code>BLOB_STORAGE_DIR</code>.</p>

<!-- 5. AI SERVICE -->
<h2>5. AI service</h2>
<p>The FastAPI service in <code>services/ai-service</code> hosts the agent registry,
orchestration logic, retrieval, financial engine, scenario engine, and report rendering.
The Anthropic SDK is used directly with <code>max_retries=5</code> and a 300-second timeout
to absorb transient 529 overloaded errors.</p>
<h3>Model routing</h3>
<table>
<tr><th>Class</th><th>Default model</th><th>Used for</th></tr>
<tr><td><code>reasoning</code></td><td>Claude Opus 4.7</td><td>Underwriting, thesis, risks, synthesis</td></tr>
<tr><td><code>default</code></td><td>Claude Sonnet 4.6</td><td>Doc classification, extraction, chat</td></tr>
<tr><td><code>fast</code></td><td>Claude Haiku 4.5</td><td>Categorization, summarization</td></tr>
</table>
<h3>Prompt caching</h3>
<p>Anthropic's prompt caching is enabled with a <em>cache primer</em> pattern.
Before the parallel agent batch fires, the orchestrator issues one tiny call
that writes the document context block to the cache. The parallel agents then
hit the cache instead of each writing redundant copies — saving roughly $0.19 per
Analyze run with a ~0.5-second latency cost. The cached prefix is the rendered
document chunks plus a shared system header; the per-agent system text sits
after the breakpoint and therefore varies safely.</p>
<h3>JSON parsing</h3>
<p>Agent outputs are validated through a layered parser (<code>app/agents/_json.py</code>):
strict <code>json.loads</code>, markdown-fence strip, smart-quote normalization,
<code>raw_decode</code> longest-prefix, then a brace-balanced recovery fallback. All
calls use <code>strict=False</code> to tolerate raw control characters that the model
sometimes emits inside string values.</p>

<!-- 6. DOC PROCESSOR -->
<h2>6. Document processor</h2>
<p>The FastAPI service in <code>services/doc-processor</code> handles ingestion:
detection, extraction, OCR fallback, semantic chunking, embedding, and indexing.
Each blob is processed once and the chunks are stored in both pgvector and ChromaDB
for hybrid retrieval (dense + BM25).</p>
<table>
<tr><th>Format</th><th>Extractor</th><th>Notes</th></tr>
<tr><td>PDF</td><td>PyMuPDF + pdfplumber</td><td>Layout-aware; tables extracted separately</td></tr>
<tr><td>Excel</td><td>openpyxl + pandas</td><td>Per-sheet, per-named-range chunks</td></tr>
<tr><td>Word</td><td>python-docx</td><td>Paragraph + heading-aware</td></tr>
<tr><td>PowerPoint</td><td>python-pptx</td><td>Per-slide chunks with speaker notes</td></tr>
<tr><td>Image</td><td>Tesseract OCR</td><td>Used as fallback when PDF text extraction fails</td></tr>
<tr><td>Email</td><td>mailparser</td><td>Headers + body + attachments unpack</td></tr>
<tr><td>ZIP</td><td>zipfile (recursive)</td><td>Children dispatch back to detect</td></tr>
</table>

<!-- 7. DATA LAYER -->
<h2>7. Data layer</h2>
<p>PostgreSQL 16 is the system of record. The <code>pgvector</code> extension adds dense
vector columns for semantic search. Prisma (in <code>packages/db</code>) is the schema and
client for the Node API; the Python services use SQLAlchemy against the same database.
The relational schema is the contract — both languages target it directly.</p>
<h3>Key entities</h3>
<table>
<tr><th>Entity</th><th>Purpose</th></tr>
<tr><td>User · Role · AuditLog</td><td>Auth, RBAC, governance trail</td></tr>
<tr><td>Client</td><td>Investor profile, mandate, preferences</td></tr>
<tr><td>Opportunity</td><td>The deal record — thesis, scores, AI verdict, financial analysis JSON</td></tr>
<tr><td>Document · DocumentChunk</td><td>Ingested artifacts with embeddings</td></tr>
<tr><td>ExtractedMetric · FinancialAssumption</td><td>Normalized financial data points</td></tr>
<tr><td>Risk · Gap</td><td>Diligence findings (category, severity, mitigation)</td></tr>
<tr><td>Scenario · ScenarioRun</td><td>What-if analyses</td></tr>
<tr><td>ChatThread · ChatMessage · Citation</td><td>Conversational intelligence with source traceability</td></tr>
<tr><td>Report</td><td>Generated memos and decks (path + payload snapshot)</td></tr>
<tr><td>LlmCall</td><td>Per-call telemetry: model, tokens, cache hits, latency, cost</td></tr>
</table>

<!-- 8. AGENT PIPELINE -->
<h2>8. AI agent pipeline</h2>
<p>The auto-Analyze pipeline that runs against an opportunity uses six agents in two stages.
The first four execute concurrently via <code>asyncio.gather(return_exceptions=True)</code>;
the orchestrator survives partial failures and returns whatever succeeded.</p>
<h3>Stage 1 — Parallel (concurrent)</h3>
<table>
<tr><th>Agent</th><th>Model</th><th>Output</th></tr>
<tr><td>Thesis Writer</td><td>Opus</td>
    <td>Thesis paragraphs, executive summary, SWOT, bull/base/bear cases, scores (opportunity / risk / confidence 0–10)</td></tr>
<tr><td>Risk Analyst</td><td>Opus</td>
    <td>List of risks: category, title, description, severity, likelihood, mitigation, citations</td></tr>
<tr><td>Gap Agent</td><td>Opus</td>
    <td>Diligence gaps + IC-readiness score (0–10)</td></tr>
<tr><td>Market Researcher</td><td>Opus + web_search tool</td>
    <td>Public-data validation of sponsor's market thesis with URL citations; TAILWIND / NEUTRAL / HEADWIND verdict</td></tr>
</table>
<h3>Stage 2 — Sequential</h3>
<table>
<tr><th>Agent</th><th>Model</th><th>Output</th></tr>
<tr><td>Financial Analyst</td><td>Opus</td>
    <td>Verdict (STRONG / MARGINAL / WEAK / INSUFFICIENT_DATA), narrative paragraphs,
        headline metrics, sensitivities, weak-assumption callouts, full metric extraction.
        Runs AFTER the parallel batch so it can interpret numbers against the risk + market findings.</td></tr>
<tr><td>Synthesis Agent</td><td>Opus</td>
    <td>Overall AI verdict (PROCEED / PROCEED_WITH_CONDITIONS / REJECT / NEED_MORE_INFO),
        rationale, top reasons for / against, critical questions, next steps, watchpoints</td></tr>
</table>
<p>Additional specialist agents exist as files in <code>app/agents/</code> (Document Analyst,
Legal Reviewer, Scenario Agent, Portfolio Exposure, Financial Model Validator, IC Challenger,
Market Stress Agent) and are dispatched by the chat orchestrator based on user intent,
but they are NOT part of the auto-Analyze pipeline.</p>
<h3>Re-run economics</h3>
<p>A dedicated <code>/analyze/opportunity/financial-only</code> endpoint reuses the prior
parallel-stage outputs from the database and re-runs only the Financial Analyst +
Synthesis Agent — useful when iterating on financial prompts. Cost: ~$1 per run vs ~$3–4
for a full Analyze; latency: ~30 seconds vs ~100 seconds.</p>

<!-- 9. END-TO-END -->
<h2>9. End-to-end flow: Analyze an opportunity</h2>
<ol>
  <li>Analyst clicks <strong>Analyze</strong> on the Understanding tab.</li>
  <li>Node API authenticates, audits, forwards to AI service <code>/analyze/opportunity</code>.</li>
  <li>AI service retrieves the top-32 chunks via hybrid search (pgvector dense + BM25)
      with reciprocal-rank fusion.</li>
  <li>Cache primer fires (one tiny call) to write the document context to Anthropic's prompt cache.</li>
  <li>Parallel agents fire (thesis, risk, gap, market). Each reads the cache, runs ~30s.</li>
  <li>Risk and Gap structured outputs are persisted to <code>Risk</code> and <code>Gap</code> tables.
      Thesis fields land on the <code>Opportunity</code> row.</li>
  <li>Financial Analyst runs with all the above as cross-agent context.</li>
  <li>Synthesis Agent runs with everything including the financial verdict.</li>
  <li>Synthesis fields and <code>financialAnalysis</code> JSON are written to the Opportunity row.</li>
  <li>Response returns to the browser; the page refreshes silently (no full-page loading state).</li>
</ol>
<p>Typical total runtime: 100–130 seconds. Typical cost: $3–4 per run with cache hits.</p>

<!-- 10. SECURITY -->
<h2>10. Security model</h2>
<ul>
  <li><strong>Auth</strong> — JWT bearer tokens, 12-hour TTL, signed with <code>JWT_SECRET</code>.</li>
  <li><strong>RBAC</strong> — Role-based access via the <code>User.role</code> column
      (ADMIN, PRINCIPAL, ANALYST, ACQUISITIONS_DIRECTOR, COMPLIANCE).</li>
  <li><strong>Audit</strong> — Every state-changing route writes an <code>AuditLog</code> entry
      with actor, action, entity, metadata, and IP.</li>
  <li><strong>Blob encryption</strong> — AES-256-GCM with a per-blob nonce. Master key in
      <code>.env</code>; envelope unwrapping happens only inside the API services.</li>
  <li><strong>API keys</strong> — Anthropic and Voyage keys live in <code>.env</code>, never
      reach the browser. The Node API proxies all Anthropic calls.</li>
  <li><strong>Data locality</strong> — All persistent state stays on the host (Docker volumes
      and local disk). Only outbound LLM calls leave the machine.</li>
</ul>

<!-- 11. BACKUP -->
<h2>11. Backup &amp; recovery</h2>
<p>Two PowerShell scripts under <code>scripts/</code> handle backup and restore:</p>
<ul>
  <li><code>backup.ps1</code> — produces a dated folder with <code>db.dump</code> (pg_dump custom-format)
      and <code>blobs.zip</code> (the <code>data/blobs</code> tree). Default retention: 14 backups.</li>
  <li><code>restore.ps1</code> — accepts a backup folder, drops the database, restores the
      dump, and unzips blobs (existing blobs are moved aside, not deleted).</li>
</ul>
<p>ChromaDB is intentionally NOT in the backup — its contents are derivable by
re-ingesting documents. The <code>.env</code> is also excluded; it should live in a password
manager.</p>

<!-- 12. DEPLOYMENT -->
<h2>12. Deployment model</h2>
<p>Stone Gate is designed local-first. The full stack runs on a developer laptop or a single
on-prem server via <code>docker compose up</code>. There is no SaaS dependency beyond the
outbound Anthropic API call. The encrypted blob store, Postgres, and Redis all live on the
host filesystem.</p>
<p>For multi-user firm deployment, the same compose file deploys cleanly behind a reverse
proxy on a single VM. Horizontal scaling is not currently a requirement — the workload is
single-firm, low-concurrency. The architecture leaves room for cloud single-tenant
deployment in the future (Phase 4 roadmap).</p>

</body></html>
"""


def render_arch_pdf(out_path: str) -> str:
    """Render the architecture PDF using the same WeasyPrint setup as IC memos."""
    env = Environment(loader=BaseLoader(), autoescape=False, trim_blocks=True, lstrip_blocks=True)
    from datetime import datetime
    html = env.from_string(ARCH_HTML).render(
        logo_data=_logo_data_uri(LOGO_PATH),
        generated_at=datetime.now().strftime("%B %Y"),
    )
    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    HTML(string=html).write_pdf(out_path)
    return out_path


# ─── Non-technical Platform Overview PPTX ────────────────────────────

def render_overview_pptx(out_path: str) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    counter = {"n": 1}

    def _next_num() -> int:
        counter["n"] += 1
        return counter["n"]

    def add_slide(title: str, bullets: list[str], subtitle: str | None = None):
        slide = prs.slides.add_slide(blank)
        _add_content_chrome(slide, prs, _next_num())
        body_top = _add_title_with_rule(slide, title, subtitle)
        body = slide.shapes.add_textbox(
            Inches(0.7), Inches(body_top), Inches(12.0),
            Inches(7.5 - body_top - 0.9),
        ).text_frame
        body.word_wrap = True
        if not bullets:
            return
        for i, b in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = f"•  {b}"
            p.font.size = Pt(15)
            p.space_after = Pt(8)

    # ── Title slide ────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    if LOGO_PATH.exists():
        slide.shapes.add_picture(
            str(LOGO_PATH),
            Inches(4.7), Inches(1.6),
            height=Inches(1.7),
        )
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.9)).text_frame
    tb.text = "Stone Gate"
    tb.paragraphs[0].font.size = Pt(44)
    tb.paragraphs[0].font.bold = True
    tb.paragraphs[0].alignment = 1
    tb.paragraphs[0].font.color.rgb = PRIMARY_RGB
    sub = tb.add_paragraph()
    sub.text = "An AI investment committee, in software form"
    sub.alignment = 1
    sub.font.size = Pt(18)
    sub.font.color.rgb = RGBColor(0x6B, 0x6B, 0x6B)
    foot = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4)).text_frame
    foot.text = "Platform overview · for principals, analysts, and decision-makers"
    foot.paragraphs[0].alignment = 1
    foot.paragraphs[0].font.size = Pt(11)
    foot.paragraphs[0].font.color.rgb = RGBColor(0x6B, 0x6B, 0x6B)

    # ── 2. The problem ────────────────────────────────────────
    add_slide(
        "The problem we solve",
        [
            "Diligence on a single real-estate deal eats weeks of analyst time — reading PDFs, extracting metrics by hand, chasing down gaps.",
            "Most of that work is repetitive: every deal needs the same shape of IC memo, the same risk taxonomy, the same scenarios.",
            "Findings end up scattered across spreadsheets and email threads. The audit trail is thin and the institutional memory is in people's heads.",
            "When a deal is reviewed at IC, contradictions or missing pieces only surface in the room — too late to do anything about them.",
        ],
    )

    # ── 3. What Stone Gate is ─────────────────────────────────
    add_slide(
        "What Stone Gate is",
        [
            "A local-first platform that takes a data room as input and produces an IC-ready recommendation as output.",
            "Built around a team of AI specialists — each one focused on a specific kind of analysis (financials, risks, market, gaps).",
            "Every claim is sourced back to the underlying document, page, and paragraph. Nothing is invented; nothing is unverifiable.",
            "Designed for a single firm. Confidential data stays on your machines — only the outbound model calls leave the network.",
        ],
    )

    # ── 4. The six-stage workflow ─────────────────────────────
    add_slide(
        "The investment workflow",
        [
            "Context — capture the briefing, the client mandate, the deal facts.",
            "Initial Screening — upload the data room; documents get classified and processed.",
            "Understanding — the AI team analyzes; outputs the thesis, financial verdict, risks, gaps, market check, and a synthesis verdict.",
            "CIO Review — the principal explores via chat, runs scenarios, stress-tests.",
            "Recommendation — generate the IC memo, executive summary, and presentation deck; record the decision.",
            "Monitoring — once approved, track performance against the underwriting.",
        ],
        subtitle="Six stages, each with a clear gate. The platform tells you where you are and what's missing.",
    )

    # ── 5. Meet the AI team ───────────────────────────────────
    add_slide(
        "Meet the AI team",
        [
            "Thesis Writer — drafts the investment story, scores opportunity / risk / confidence.",
            "Risk Analyst — enumerates risks across 11 categories, severity-ranked with mitigations.",
            "Gap Agent — flags missing diligence items and scores IC-readiness.",
            "Market Researcher — independently validates the sponsor's market thesis against public data, with web search.",
            "Financial Analyst — pressure-tests the numbers; produces a verdict, sensitivities, weak-assumption callouts.",
            "Synthesis Agent — reads everyone's findings; produces the final AI recommendation and rationale.",
        ],
        subtitle="Six specialists. Each one has a single job and a calibrated rubric.",
    )

    # ── 6. How an Analyze works ───────────────────────────────
    add_slide(
        "How an Analyze run works",
        [
            "The first four agents work in parallel — like four analysts assigned to the same deal at the same time.",
            "Once they're done, the Financial Analyst takes over — it sees what the others found and pressure-tests the numbers against the risks and market findings.",
            "Finally, the Synthesis Agent reads everything and writes the institutional verdict — PROCEED, REJECT, or NEED MORE INFO — with reasons for and against.",
            "Total time: roughly 100 seconds. Total cost: about $3–4 in model calls.",
            "If you only want to update the financial view (e.g., after tightening the prompt), there's a 30-second \"financials-only\" re-run that costs about $1.",
        ],
    )

    # ── 7. What you actually see ──────────────────────────────
    add_slide(
        "What the analyst sees",
        [
            "A single Understanding page with collapsible sections: Scores → AI Recommendation → Thesis → Market Research → Financials → Risks → Gaps → Reports.",
            "Every claim links back to the source document with a citation chip (chunk + page).",
            "The Financials section is itself a mini-memo: verdict pill, narrative paragraphs, headline metrics, sensitivities, weak-assumption callouts.",
            "The Risks table is filterable by Category and Severity; a heatmap above shows the distribution at a glance.",
            "Filters are instant. No round-trips. Everything is keyboard-friendly.",
        ],
    )

    # ── 8. Reports ────────────────────────────────────────────
    add_slide(
        "Reports & exports",
        [
            "IC Memorandum — a full 12–15 page institutional memo (PDF), with executive summary, thesis, financial analysis, risks, market analysis, recommendation.",
            "Executive Summary — a tight 2–3 page principal brief (PDF), focused on the verdict and the three things that matter most.",
            "IC Presentation Deck — a 14-slide deck (PowerPoint) covering economics, financial verdict, sensitivities, top risks, bull/base/bear, AI verdict, next steps.",
            "All three are generated from the same underlying analysis, so they're always internally consistent.",
            "The same Stone Gate brand on all three — maroon header band, white body, sourced citations.",
        ],
    )

    # ── 9. Chat & insight ─────────────────────────────────────
    add_slide(
        "Chat with the deal",
        [
            "An Insights chat panel on every opportunity — ask anything in plain English: \"what are the top three risks?\", \"what happens if the exit cap widens to 6.5%?\".",
            "The chat is grounded in the documents you uploaded. Every answer cites the source.",
            "The system routes the question to the right specialist — risk questions go to the Risk Analyst, financial questions go to the Financial Analyst, and so on.",
            "Long threads get summarized and folded back into the opportunity's memory, so you don't have to re-explain context next week.",
        ],
    )

    # ── 10. Documents ─────────────────────────────────────────
    add_slide(
        "Document intelligence",
        [
            "Upload PDFs, Excel models, decks, Word docs, emails, even ZIPs — all handled.",
            "Each document is classified (\"financial model\", \"market study\", \"legal\", etc.) and broken into searchable chunks.",
            "Scanned PDFs get OCR automatically; image-only files are extracted as best they can be.",
            "Every chunk is embedded and indexed, so retrieval is semantic — \"capital structure\" finds the right paragraph even if the deck called it \"funding stack\".",
        ],
    )

    # ── 11. Security ──────────────────────────────────────────
    add_slide(
        "Security & privacy",
        [
            "Local-first: everything persists on your hardware. There is no SaaS backend.",
            "Document blobs encrypted at rest with AES-256-GCM, per-blob nonces, master key in the local .env.",
            "Role-based access — Principal, Analyst, Director of Acquisitions, Compliance, Admin — with a full audit trail on every action.",
            "Only outbound traffic: anonymous API calls to Anthropic (with prompt caching to minimize what gets sent and to control cost).",
            "Backup and restore scripts ship in the repo. Daily Postgres dumps + blob archives, retention configurable.",
        ],
    )

    # ── 12. Where we are ──────────────────────────────────────
    add_slide(
        "Where we are today",
        [
            "Foundation, MVP and V1 features are live: ingestion, multi-agent analysis, financial deep-read, IC memo / exec summary / deck generation, chat with citations.",
            "Telemetry on every model call — token counts, cache hits, cost, latency — so you can see exactly what each Analyze cost.",
            "Backup tooling and INSTALL documentation ready for deployment on a second machine.",
            "Code is version-controlled, off-machine, and recoverable in under 30 minutes on a clean machine.",
        ],
    )

    # ── 13. What's next ───────────────────────────────────────
    add_slide(
        "What's next",
        [
            "Portfolio overlap & exposure — surface conflicts and concentration before they become problems.",
            "Monte Carlo and cognitive-bias warnings — quantify uncertainty in the underwriting.",
            "Prompt versioning and A/B harness — measure analyst-grade improvements over time.",
            "CRM and PMS integrations — bidirectional sync with Salesforce, HubSpot, and your portfolio management system.",
            "Cloud single-tenant deployment for firms that prefer it over local-first.",
        ],
    )

    # ── 14. Closing ───────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    _add_content_chrome(slide, prs, _next_num())
    _add_title_with_rule(slide, "Stone Gate")
    body = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(12.0), Inches(4)).text_frame
    body.text = "Faster diligence. Defensible decisions. Institutional memory."
    body.paragraphs[0].font.size = Pt(28)
    body.paragraphs[0].font.bold = True
    body.paragraphs[0].font.color.rgb = PRIMARY_RGB
    sub = body.add_paragraph()
    sub.text = "Ready to use today, designed to keep improving."
    sub.font.size = Pt(16)
    sub.font.color.rgb = RGBColor(0x6B, 0x6B, 0x6B)
    sub.space_before = Pt(10)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def main() -> None:
    out_root = "/app/data/blobs/platform-docs"
    pdf_path = os.path.join(out_root, "stone-gate-architecture.pdf")
    pptx_path = os.path.join(out_root, "stone-gate-overview.pptx")
    print(f"Generating architecture PDF -> {pdf_path}")
    render_arch_pdf(pdf_path)
    print(f"Generating overview PPTX    -> {pptx_path}")
    render_overview_pptx(pptx_path)
    print("Done.")


if __name__ == "__main__":
    main()
