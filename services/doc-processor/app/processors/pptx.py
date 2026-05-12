"""PPTX processor — extracts slide text + speaker notes per slide."""
from __future__ import annotations

import io
from dataclasses import dataclass

from pptx import Presentation


@dataclass
class SlideContent:
    page: int  # 1-indexed slide number
    text: str


def extract_pptx(buf: bytes) -> tuple[list[SlideContent], int]:
    prs = Presentation(io.BytesIO(buf))
    slides: list[SlideContent] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    line = "".join(r.text for r in p.runs).strip()
                    if line:
                        parts.append(line)
            # Tables
            if getattr(shape, "has_table", False):
                tbl = shape.table
                for row in tbl.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append("\t".join(cells))
        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Speaker notes] {notes}")

        text = f"Slide {i}\n" + "\n".join(parts) if parts else f"Slide {i}\n(empty)"
        slides.append(SlideContent(page=i, text=text))
    return slides, len(slides)
